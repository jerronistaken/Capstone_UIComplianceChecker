import os
import re
from io import BytesIO

import pandas as pd
import streamlit as st
import textstat

try:
    import google.generativeai as genai
except ImportError:
    genai = None

from PIL import Image
try:
    import pytesseract
    tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.exists(tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
except ImportError:
    pytesseract = None
except Exception:
    pytesseract = None

from docx import Document
from pypdf import PdfReader
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# =========================
# CONFIG (inline)
# =========================

BRAND_KEYWORDS = ["JTC"]   # ← Add more brand-specific keywords as needed with "," e.g. "JTC", "JTC Group", "JTC Corporation"
FORBIDDEN_WORDS = ["TBD", "Lorem ipsum", "placeholder", "TODO"]
MIN_WORD_COUNT = 100
MAX_SENTENCE_LENGTH = 25
MIN_READABILITY_SCORE = 30

# =========================
# PAGE CONFIG
# =========================

st.set_page_config(page_title="Document Compliance Checker", layout="wide")
st.title("Automated Document Compliance & UX Checker")
st.write("Upload PDF, DOCX, or PPTX files for automated analysis. Set GOOGLE_API_KEY to enable Gemini suggestions.")

# =========================
# TEXT EXTRACTION
# =========================

def ocr_image_to_text(image):
    if pytesseract is None:
        return ""
    try:
        return pytesseract.image_to_string(image).strip()
    except Exception:
        return ""


def extract_images_from_docx(doc):
    images = []
    for rel in doc.part.rels.values():
        reltype = getattr(rel, "reltype", "")
        if "image" not in str(reltype).lower():
            continue
        target_part = getattr(rel, "target_part", None)
        blob = getattr(target_part, "blob", None)
        if not blob:
            continue
        try:
            images.append(Image.open(BytesIO(blob)))
        except Exception:
            continue
    return images


def extract_text_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])

    ocr_texts = []
    for image in extract_images_from_docx(doc):
        image_text = ocr_image_to_text(image)
        if image_text:
            ocr_texts.append(image_text)

    if ocr_texts:
        text += "\n" + "\n".join(ocr_texts)
    return text


def extract_text_from_pdf(file):
    reader = PdfReader(file)
    text_chunks = []
    ocr_texts = []

    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text_chunks.append(extracted)

        for image_file in getattr(page, "images", []):
            try:
                image = image_file.image if getattr(image_file, "image", None) is not None else Image.open(BytesIO(image_file.data))
            except Exception:
                continue
            image_text = ocr_image_to_text(image)
            if image_text:
                ocr_texts.append(image_text)

    if ocr_texts:
        text_chunks.append("\n".join(ocr_texts))
    return "\n".join(text_chunks)


def extract_images_from_pptx(presentation):
    images = []

    def collect_images(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                images.append(Image.open(BytesIO(shape.image.blob)))
            except Exception:
                pass
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                collect_images(child)
        elif hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for cell_shape in cell.shapes:
                        collect_images(cell_shape)

    for slide in presentation.slides:
        for shape in slide.shapes:
            collect_images(shape)

    return images


def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text_chunks = []
    ocr_texts = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text:
                text_chunks.append(shape.text)
            elif hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_chunks.append(cell.text)

    for image in extract_images_from_pptx(presentation):
        image_text = ocr_image_to_text(image)
        if image_text:
            ocr_texts.append(image_text)

    if ocr_texts:
        text_chunks.append("\n".join(ocr_texts))
    return "\n".join(text_chunks)


def extract_docx_section_hints(file):
    doc = Document(file)
    sections = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = getattr(getattr(para, "style", None), "name", "")
        if "Heading" in style_name or text.isupper() or text.endswith(":"):
            sections.append(text if len(text) <= 80 else text[:77] + "...")
        if len(sections) >= 5:
            break

    if not sections:
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                sections.append(text if len(text) <= 80 else text[:77] + "...")
            if len(sections) >= 3:
                break

    if not sections:
        return "Document has no obvious headings or section markers."

    return "Sections: " + " | ".join(sections)


def extract_pdf_page_hints(file):
    reader = PdfReader(file)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        extracted = page.extract_text()
        if not extracted:
            pages.append(f"Page {i}: no extractable text")
        else:
            snippet = " ".join(extracted.strip().split())
            pages.append(f"Page {i}: {snippet[:100]}{'...' if len(snippet) > 100 else ''}")
        if len(pages) >= 5:
            break

    return " | ".join(pages)


def extract_pptx_slide_hints(file):
    presentation = Presentation(file)
    slides = []
    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text:
                slide_text.append(shape.text.strip().replace("\n", " "))
        if slide_text:
            snippet = " ".join(slide_text)[:120]
            slides.append(f"Slide {i}: {snippet}{'...' if len(slide_text) > 0 and len(snippet) == 120 else ''}")
        else:
            slides.append(f"Slide {i}: no visible text; may contain images")
        if len(slides) >= 5:
            break

    return " | ".join(slides)


def check_bullet_formatting(text):
    lines = text.splitlines()
    bullet_issues = []
    indent_levels = []

    for idx, line in enumerate(lines, 1):
        match = re.match(r'^(\s*)([-*•+])(\s*)(.*)$', line)
        if match and match.group(4).strip():
            indent = len(match.group(1))
            spacing = match.group(3)
            if spacing != " ":
                bullet_issues.append(f"Line {idx}: bullet marker should be followed by a single space.")
            indent_levels.append(indent)

    indent_issue = None
    unique_indents = sorted(set(indent_levels))
    if len(unique_indents) > 1:
        diffs = [unique_indents[i] - unique_indents[i - 1] for i in range(1, len(unique_indents))]
        if any(indent % 2 != 0 for indent in unique_indents):
            indent_issue = "Bullet indentation should use even-numbered spaces."
        elif len(set(diffs)) > 1:
            indent_issue = "Bullet indentation uses inconsistent nesting increments."

    if indent_issue:
        bullet_issues.append(indent_issue)

    return bullet_issues


def generate_ai_suggestions(text, results):
    if genai is None:
        return ["Install the google-generativeai package and add it to requirements for Gemini suggestions."]

    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        return ["Set the GOOGLE_API_KEY environment variable to enable Gemini suggestions."]

    try:
        genai.configure(api_key=api_key)
        summary = (
            "Review the following document text and the detected issues, then provide concise improvement suggestions. "
            "Focus on spacing, bullet formatting, indentation, grammar, readability, and placeholder content."
        )
        prompt_text = (
            f"Document text:\n{text[:2000]}\n\nDetected issues:\n"
            f"Missing keywords: {results['Missing Brand Keywords']}\n"
            f"Forbidden words: {results['Forbidden Words Found']}\n"
            f"Grammar issues: {results['Grammar Issues']}\n"
            f"Readability score: {results['Readability Score']}\n"
            f"Average sentence length: {results['Average Sentence Length']}\n"
            f"Bullet formatting issues: {results.get('Bullet Formatting Details')}\n\n"
            f"Document context: {results.get('Document Context', 'No context available')}\n"
            "If this is a PPTX file, mention the slide number for each recommendation. "
            "If this is a DOCX or PDF file, mention the page or section where the issue occurs."
        )

        model = os.getenv("GOOGLE_GEMINI_MODEL", "gemini-flash-latest")
        try:
            # The installed `google.generativeai` package exposes a
            # `GenerativeModel` -> `start_chat()` -> `send_message()` API.
            # Normalize model name if the environment variable doesn't
            # include the `models/` prefix.
            model_name = model if str(model).startswith("models/") else f"models/{model}"
            gm = genai.GenerativeModel(model_name)
            chat = gm.start_chat()
            response = chat.send_message(summary + "\n\n" + prompt_text)
            # Response objects on this client expose `text`.
            suggestion = getattr(response, 'text', str(response))
            return [suggestion]
        except Exception as e:
            return [f"Gemini suggestion unavailable: {str(e)}"]
    except Exception as e:
        return [f"Gemini suggestion unavailable: {str(e)}"]

# =========================
# ANALYSIS FUNCTIONS
# =========================

def check_brand_keywords(text):
    return [kw for kw in BRAND_KEYWORDS if kw.lower() not in text.lower()]


def check_forbidden_words(text):
    return [w for w in FORBIDDEN_WORDS if w.lower() in text.lower()]


def check_grammar(text):
    """
    Lightweight grammar heuristic — counts common grammar signals
    without requiring Java or heavy external services.
    Flags: double spaces, repeated words, missing space after punctuation,
    and i not capitalised.
    """
    issues = 0
    issues += len(re.findall(r'  +', text))                        # double spaces
    issues += len(re.findall(r'\b(\w+)\s+\1\b', text, re.I))      # repeated words
    issues += len(re.findall(r'[.!?,;][a-zA-Z]', text))           # missing space after punctuation
    issues += len(re.findall(r'\bi\b', text))                      # uncapitalised "i"
    return issues


def readability_analysis(text):
    return round(textstat.flesch_reading_ease(text), 2)


def average_sentence_length(text):
    sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip()]
    if not sentences:
        return 0
    total_words = sum(len(s.split()) for s in sentences)
    return round(total_words / len(sentences), 2)


def generate_feedback(results):
    feedback = []

    if results['Missing Brand Keywords']:
        feedback.append(f"Missing branding keywords: {', '.join(results['Missing Brand Keywords'])}.")

    if results['Forbidden Words Found']:
        feedback.append(f"Forbidden placeholder terms detected: {', '.join(results['Forbidden Words Found'])}.")

    if results['Bullet Formatting Issues'] > 0:
        feedback.append("Bullet formatting issues detected. Review bullet spacing and indentation.")

    if results['Grammar Issues'] > 10:
        feedback.append("High number of grammar issues detected.")

    if results['Readability Score'] < MIN_READABILITY_SCORE:
        feedback.append("Document readability is too difficult for average readers.")

    if results['Average Sentence Length'] > MAX_SENTENCE_LENGTH:
        feedback.append("Sentences are excessively long and may reduce readability.")

    if results['Word Count'] < MIN_WORD_COUNT:
        feedback.append("Document content may be insufficient or incomplete.")

    if not feedback:
        feedback.append("Document passed all major checks.")

    return feedback


def calculate_score(results):
    score = 100
    score -= len(results['Missing Brand Keywords']) * 10
    score -= len(results['Forbidden Words Found']) * 10
    score -= min(results['Grammar Issues'], 20)
    score -= min(results['Bullet Formatting Issues'] * 5, 20)

    if results['Readability Score'] < MIN_READABILITY_SCORE:
        score -= 15

    if results['Average Sentence Length'] > MAX_SENTENCE_LENGTH:
        score -= 10

    if results['Word Count'] < MIN_WORD_COUNT:
        score -= 15

    return max(score, 0)


def determine_status(score):
    if score >= 85:
        return "PASS"
    elif score >= 60:
        return "WARNING"
    else:
        return "FAIL"

# =========================
# EXCEL EXPORT
# =========================

def export_to_excel(results_list, output_path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Compliance Report"

    headers = [
        "Document Name", "Status", "Compliance Score", "Word Count",
        "Grammar Issues", "Readability Score", "Avg Sentence Length",
        "Bullet Formatting Issues", "Missing Brand Keywords", "Forbidden Words", "Feedback", "AI Suggestions"
    ]
    ws.append(headers)

    header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = header_fill

    for result in results_list:
        ws.append([
            result['Document Name'],
            result['Status'],
            result['Compliance Score'],
            result['Word Count'],
            result['Grammar Issues'],
            result['Readability Score'],
            result['Average Sentence Length'],
            result['Bullet Formatting Issues'],
            ", ".join(result['Missing Brand Keywords']),
            ", ".join(result['Forbidden Words Found']),
            " | ".join(result['Feedback']),
            " | ".join(result.get('AI Suggestions', []))
        ])

    for col in ws.columns:
        max_length = max((len(str(cell.value)) for cell in col if cell.value), default=10)
        ws.column_dimensions[col[0].column_letter].width = max_length + 5

    wb.save(output_path)

# =========================
# MAIN UI
# =========================

uploaded_files = st.file_uploader(
    "Upload Documents",
    type=['pdf', 'docx', 'pptx'],
    accept_multiple_files=True
)

if uploaded_files:
    results_list = []

    for uploaded_file in uploaded_files:
        filename = uploaded_file.name
        extension = filename.split('.')[-1].lower()

        try:
            if extension == 'docx':
                text = extract_text_from_docx(uploaded_file)
                uploaded_file.seek(0)
                location_hints = extract_docx_section_hints(uploaded_file)
            elif extension == 'pdf':
                text = extract_text_from_pdf(uploaded_file)
                uploaded_file.seek(0)
                location_hints = extract_pdf_page_hints(uploaded_file)
            elif extension == 'pptx':
                text = extract_text_from_pptx(uploaded_file)
                uploaded_file.seek(0)
                location_hints = extract_pptx_slide_hints(uploaded_file)
            else:
                st.warning(f"Unsupported file: {filename}")
                continue

            bullet_issues = check_bullet_formatting(text)
            results = {
                'Document Context': location_hints,
                'Document Name': filename,
                'Word Count': len(text.split()),
                'Missing Brand Keywords': check_brand_keywords(text),
                'Forbidden Words Found': check_forbidden_words(text),
                'Grammar Issues': check_grammar(text),
                'Readability Score': readability_analysis(text),
                'Average Sentence Length': average_sentence_length(text),
                'Bullet Formatting Issues': len(bullet_issues),
                'Bullet Formatting Details': bullet_issues
            }

            results['Feedback'] = generate_feedback(results)
            results['Compliance Score'] = calculate_score(results)
            results['Status'] = determine_status(results['Compliance Score'])
            results['AI Suggestions'] = generate_ai_suggestions(text, results)
            results_list.append(results)

        except Exception as e:
            st.error(f"Error processing {filename}: {str(e)}")

    if results_list:
        st.subheader("Analysis Results")
        st.dataframe(pd.DataFrame(results_list))

        for result in results_list:
            st.markdown("---")
            st.subheader(result['Document Name'])

            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Compliance Score", result['Compliance Score'])
            with col2:
                st.metric("Status", result['Status'])
            with col3:
                st.metric("Grammar Issues", result['Grammar Issues'])

            st.write("### Feedback")
            for fb in result['Feedback']:
                st.warning(fb)

            st.write("### AI Suggestions")
            for suggestion in result['AI Suggestions']:
                st.info(suggestion)

        os.makedirs('reports', exist_ok=True)
        excel_path = 'reports/compliance_report.xlsx'
        export_to_excel(results_list, excel_path)

        with open(excel_path, 'rb') as f:
            st.download_button(
                label="Download Excel Report",
                data=f,
                file_name="compliance_report.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )